from dotenv import load_dotenv
import openai
from pptx import Presentation
from gtts import gTTS
import json
import os
import re
from azure.cognitiveservices.speech import SpeechConfig, SpeechSynthesisOutputFormat, SpeechSynthesizer

# Load environment variables
load_dotenv()

# API Configuration
try:
    DEEPSEEK_API_KEY = os.environ['DEEPSEEK_API_KEY']
    AZURE_SPEECH_KEY = os.environ['AZURE_SPEECH_KEY']
    AZURE_SPEECH_REGION = os.environ['AZURE_SPEECH_REGION']
    DEEPSEEK_API_URL = "https://api.deepseek.com/v1"
except KeyError:
    raise KeyError("API KEY not found. Make sure it's set in your .env file or environment variables.")


def generate_ppt_structure(user_input):
    """Generate structured PowerPoint content using LLM API"""
    from openai import OpenAI
    
    # Initialize client with base URL and API key
    client = OpenAI(
        api_key=DEEPSEEK_API_KEY,
        base_url=DEEPSEEK_API_URL
    )

    prompt = f"""
作为专业PPT设计师，请根据以下需求生成详细结构，保证所有内容为正式英文输出：
1. 包含5-8页幻灯片
2. 每页必须包含：
   - 标题（不超过10个单词）
   - 最重要的3-5个详细要点（每个要点包含子要点进行深入解释，但要保证解释的正确性和真实性。）
   - 演讲脚本（尽量使用口语化表达，包含过渡语句，你的演讲要足够使人信服，需要你有高超的演讲技巧。）
3. 按此JSON格式返回：
{{
  "title": "主标题",
  "slides": [
    {{
      "title": "页面标题",
      "content": [
        "主要观点1（详细解释）",
        "主要观点2（详细解释）",
        ...
      ],
      "speech":"使用英文脚本，包含自然过渡词的讲解脚本 （请保持高超的演讲技巧。要讲解的清晰有趣且易于理解，避免过于专业的词语。）"
    }}
  ]
}}

用户输入：{user_input}
确保直接返回严格JSON格式，不要任何额外文本！JSON字段中避免使用换行符和特殊控制字符。
"""

    try:
        response = client.chat.completions.create(
            model="deepseek-reasoner",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4
        )
        
        # Parse the response
        result = response.choices[0].message.content
        
        try:
            return json.loads(result)
        except json.JSONDecodeError:
            # Clean JSON string if direct parsing fails
            sanitized = result.replace('\n', ' ').replace('\r', ' ')
            
            # Extract JSON content with regex
            json_match = re.search(r'({[\s\S]*})', sanitized)
            if json_match:
                try:
                    return json.loads(json_match.group(1))
                except:
                    pass
            
            # Return a simple structure if all parsing attempts fail
            return {
                "title": user_input,
                "slides": [
                    {
                        "title": "Introduction",
                        "content": ["Could not parse AI response. This is a fallback slide."],
                        "speech": "I'll be discussing " + user_input + " today."
                    }
                ]
            }
    except Exception as e:
        print(f"API call failed: {str(e)}")
        raise

def create_ppt(data, filename="output.pptx"):
    """Generate PowerPoint file from structured data"""
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = data['title']
    
    # Content slides
    for slide_data in data['slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data['title']
        content.text = '\n'.join(slide_data['content'])
    
    prs.save(filename)
    return filename

def generate_speeches(data, output_dir="audio"):
    """Generate speech audio for each slide"""
    os.makedirs(output_dir, exist_ok=True)
    
    # Configure Azure Speech Service for more natural male voice
    speech_config = SpeechConfig(subscription=AZURE_SPEECH_KEY, region=AZURE_SPEECH_REGION)
    speech_config.set_speech_synthesis_output_format(
        SpeechSynthesisOutputFormat.Audio24Khz160KBitRateMonoMp3)  # High quality audio
    
    # Choose an engaging female voice for English presentations
    speech_config.speech_synthesis_voice_name = "en-US-JennyNeural"
    
    synthesizer = SpeechSynthesizer(speech_config=speech_config)
    
    audio_files = []
    for i, slide in enumerate(data['slides']):
        # Use SSML with prosody adjustments for more engaging, enthusiastic style
        ssml = f"""
        <speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis"
               xmlns:mstts="https://www.w3.org/2001/mstts" xml:lang="en-US">
            <voice name="{speech_config.speech_synthesis_voice_name}">
                <prosody rate="1.0" pitch="+5%">
                    <mstts:express-as style="cheerful" styledegree="1.2">
                        {slide['speech']}
                    </mstts:express-as>
                </prosody>
            </voice>
        </speak>
        """
        
        try:
            # Try Azure Speech Service first
            result = synthesizer.speak_ssml_async(ssml).get()
            filename = f"{output_dir}/slide_{i+1}.mp3"
            with open(filename, "wb") as audio_file:
                audio_file.write(result.audio_data)
            audio_files.append(filename)
        except Exception as e:
            print(f"Azure Speech failed for slide {i+1}: {str(e)}")
            # Fall back to gTTS
            try:
                tts = gTTS(text=slide['speech'], lang='zh-cn', slow=False)
                filename = f"{output_dir}/slide_{i+1}.mp3"
                tts.save(filename)
                audio_files.append(filename)
            except Exception as e2:
                print(f"Speech generation failed: {str(e2)}")
    
    return audio_files

if __name__ == "__main__":
    # Simple CLI interface for testing
    user_input = input("Enter presentation topic: ")
    ppt_data = generate_ppt_structure(user_input)
    ppt_file = create_ppt(ppt_data)
    audio_files = generate_speeches(ppt_data)
    print(f"PowerPoint saved as: {ppt_file}")
    print(f"Audio files generated: {len(audio_files)}")