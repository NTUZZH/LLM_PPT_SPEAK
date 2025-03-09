from flask import Flask, render_template, request, jsonify, send_from_directory
import os
import time
from PIL import Image, ImageDraw, ImageFont
import io
from pptx import Presentation
from MVP import generate_ppt_structure, create_ppt, generate_speeches
import json
import re

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'static/generated'
app.config['SLIDES_FOLDER'] = 'static/slides'
app.config['AUDIO_FOLDER'] = 'static/audio'

# Ensure directories exist
for folder in [app.config['UPLOAD_FOLDER'], app.config['SLIDES_FOLDER'], app.config['AUDIO_FOLDER']]:
    os.makedirs(folder, exist_ok=True)

def convert_ppt_to_images(pptx_path):
    """Convert PowerPoint slides to images"""
    # Create unique subfolder for this presentation
    timestamp = int(time.time())
    output_dir = os.path.join(app.config['SLIDES_FOLDER'], f"pres_{timestamp}")
    os.makedirs(output_dir, exist_ok=True)
    
    # First, check if the PPTX file exists
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
    
    # Get the base name of the PowerPoint file (without extension)
    ppt_basename = os.path.splitext(os.path.basename(pptx_path))[0]
    
    # LibreOffice will create a PDF with the same base name as the input file
    expected_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{ppt_basename}.pdf")
    
    # Run LibreOffice conversion
    cmd = f"soffice --headless --convert-to pdf --outdir {app.config['UPLOAD_FOLDER']} {pptx_path}"
    print(f"Running: {cmd}")
    exit_code = os.system(cmd)
    
    if exit_code != 0 or not os.path.exists(expected_pdf_path):
        print("PDF conversion failed, falling back to direct image creation")
        return create_placeholder_images(pptx_path, output_dir, timestamp)
    
    # Now convert PDF to images
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(expected_pdf_path)
        
        slide_paths = []
        for i, image in enumerate(images):
            file_path = os.path.join(output_dir, f"slide_{i+1}.jpg")
            image.save(file_path, "JPEG")
            slide_paths.append(f"/static/slides/pres_{timestamp}/slide_{i+1}.jpg")
        
        # Clean up PDF
        os.remove(expected_pdf_path)
        return slide_paths
    except Exception as e:
        print(f"Error in PDF to image conversion: {str(e)}")
        return create_placeholder_images(pptx_path, output_dir, timestamp)

# Helper function for creating placeholder images
def create_placeholder_images(pptx_path, output_dir, timestamp):
    """Create basic placeholder images with slide titles"""
    prs = Presentation(pptx_path)
    slide_paths = []
    
    for i, slide in enumerate(prs.slides):
        # Create a blank image
        img = Image.new('RGB', (960, 720), (255, 255, 255))
        # Add slide number text
        draw = ImageDraw.Draw(img)
        
        # Try to extract slide title
        title_text = f"Slide {i+1}"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                title_text = shape.text.strip() or title_text
                break
                
        # Draw text in center with proper centering
        try:
            # Try to load a system font
            font = ImageFont.truetype("Arial", 30)
        except:
            # Fall back to default
            font = ImageFont.load_default()

        # Calculate text position to center it
        text_width = draw.textlength(title_text, font=font)
        text_position = ((960 - text_width) // 2, 360)

        # Draw the text
        draw.text(text_position, title_text, fill=(0, 0, 0), font=font)
        
        file_path = os.path.join(output_dir, f"slide_{i+1}.jpg")
        img.save(file_path, "JPEG")
        slide_paths.append(f"/static/slides/pres_{timestamp}/slide_{i+1}.jpg")
    
    return slide_paths

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    user_input = request.form['topic']
    
    try:
        # Generate structured data
        raw_data = generate_ppt_structure(user_input)
        
        # If it's a string (not yet parsed JSON), parse it properly
        if isinstance(raw_data, str):
            import json
            import re
            
            # Step 1: First try direct parsing
            try:
                ppt_data = json.loads(raw_data)
            except json.JSONDecodeError as e:
                print(f"JSON Parse Error: {str(e)}")
                print(f"Error at position {e.pos}: {raw_data[max(0, e.pos-10):e.pos]}|{raw_data[e.pos:min(len(raw_data), e.pos+10)]}")
                
                # Step 2: Try to fix common JSON issues
                # Remove any potential code block markers from LLM responses
                sanitized = re.sub(r'```json|```', '', raw_data)
                sanitized = sanitized.strip()
                
                # Fix newlines in strings
                sanitized = sanitized.replace('\n', ' ').replace('\r', '')
                
                try:
                    ppt_data = json.loads(sanitized)
                except json.JSONDecodeError:
                    # Step 3: More aggressive cleaning
                    # Extract just the part that looks like valid JSON
                    json_match = re.search(r'({[\s\S]*})', sanitized)
                    if json_match:
                        try:
                            ppt_data = json.loads(json_match.group(1))
                        except:
                            # Final fallback: create a simple valid structure
                            ppt_data = {
                                "title": user_input,
                                "slides": [
                                    {
                                        "title": "Error Processing Content",
                                        "content": ["Could not parse AI response"],
                                        "speech": "Sorry, there was an error processing the content."
                                    }
                                ]
                            }
                    else:
                        raise ValueError("Could not extract valid JSON from response")
        else:
            ppt_data = raw_data
        
        # Continue with the rest of your function...
        # Create a unique filename
        timestamp = int(time.time())
        ppt_filename = f"presentation_{timestamp}.pptx"
        ppt_path = os.path.join(app.config['UPLOAD_FOLDER'], ppt_filename)
        
        # Create PPT
        create_ppt(ppt_data, ppt_path)
        
        # Generate audio files
        audio_dir = os.path.join(app.config['AUDIO_FOLDER'], f"pres_{timestamp}")
        os.makedirs(audio_dir, exist_ok=True)
        
        # Adjust generate_speeches to use our audio directory
        audio_files = generate_speeches(ppt_data, audio_dir)
        
        # Convert slide relative paths to web paths
        audio_web_paths = [f"/static/audio/pres_{timestamp}/{os.path.basename(f)}" for f in audio_files]
        
        # Convert PPT to images
        slide_paths = convert_ppt_to_images(ppt_path)
        
        # Create presentation data structure
        presentation_data = {
            'title': ppt_data['title'],
            'slides': [],
        }
        
        # First slide has no audio
        presentation_data['slides'].append({
            'image': slide_paths[0] if slide_paths else "",
            'audio': "",  # No audio for the first slide
            'title': ppt_data['slides'][0]['title'] if ppt_data['slides'] else "Introduction"
        })
        
        # Match remaining slides with audio (offset by 1)
        for i in range(1, len(slide_paths)):
            audio_index = i - 1  # Audio is offset by 1
            presentation_data['slides'].append({
                'image': slide_paths[i],
                'audio': audio_web_paths[audio_index] if audio_index < len(audio_web_paths) else "",
                'title': ppt_data['slides'][i]['title'] if i < len(ppt_data['slides']) else f"Slide {i+1}"
            })
        
        return render_template('presentation.html', presentation=presentation_data)
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    # Use environment variable for port with a default of 8080 instead of 5000
    port = int(os.environ.get('PORT', 8080))
    # In production, set debug to False
    debug = os.environ.get('FLASK_DEBUG', 'False').lower() == 'true'
    app.run(host='0.0.0.0', port=port, debug=debug)
